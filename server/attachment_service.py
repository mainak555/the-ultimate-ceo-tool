"""Attachment upload, extraction, binding, and retrieval helpers.

Extraction strategy
-------------------
Text extraction is **lazy**: nothing is extracted at upload time so the user
gets an instant response and no CPU is wasted on files that are never sent.

When the first agent run (or any resume) calls
``build_attachment_context_block``, each non-image attachment's text is
loaded as follows:

1. Check Redis ``{ns}:attachment:{session_id}:{attachment_id}:text``
   (TTL = ``REDIS_ATTACHMENT_TTL_SECONDS``, default 24 h).
2. Cache **HIT** → return text immediately (< 1 ms).
3. Cache **MISS** → download raw bytes from Azure Blob, extract, store in
   Redis, return text.  This one-time extraction is done inside
   ``asyncio.to_thread`` so the event loop is never blocked.

Images are **never** stored in Redis.  Each run/resume downloads image bytes
directly from blob and wraps them in ``autogen_core.Image`` objects for
``MultiModalMessage``.  Vision models receive actual pixel data.

MongoDB stores only metadata — no file content ever lands in the database.
"""

from __future__ import annotations

import csv
import io
import json
import logging
import re
from typing import Iterable
from uuid import uuid4

from bson import ObjectId
from bson.errors import InvalidId
from django.conf import settings
from django.core.files.uploadedfile import UploadedFile

from core.tracing import traced_function

from .db import get_collection, ATTACHMENTS_COLLECTION
from .storage_backends import build_storage_strategy
from . import util

logger = logging.getLogger(__name__)

_MAX_ATTACHMENTS_PER_MESSAGE = 10
_MAX_ATTACHMENT_BYTES = 20 * 1024 * 1024

# ---------------------------------------------------------------------------
# Redis cache helpers
# ---------------------------------------------------------------------------

def _attachment_cache_ttl() -> int:
    """Return the attachment text cache TTL (minimum 1 h)."""
    raw = int(getattr(settings, "REDIS_ATTACHMENT_TTL_SECONDS", 86400) or 86400)
    return max(3600, raw)


def _redis_namespace() -> str:
    ns = (getattr(settings, "REDIS_NAMESPACE", "product_discovery") or "product_discovery").strip()
    return ns or "product_discovery"


def _att_text_key(session_id: str, attachment_id: str) -> str:
    return f"{_redis_namespace()}:attachment:{session_id}:{attachment_id}:text"


def _att_index_key(session_id: str) -> str:
    """Redis SET key that tracks all attachment_ids cached for a session."""
    return f"{_redis_namespace()}:attachment:{session_id}:index"


def _get_redis():
    """Return the shared Redis client (reuses agent session_coordination pool).

    Returns ``None`` and logs a warning when Redis is unavailable so callers
    can fall back to blob-only extraction without crashing.
    """
    try:
        from agents.session_coordination import get_redis_client
        return get_redis_client()
    except Exception:
        logger.warning("attachments.redis_unavailable", exc_info=False)
        return None


def _redis_get_text(session_id: str, attachment_id: str) -> str | None:
    """Return cached extracted text or ``None`` on miss/error."""
    try:
        r = _get_redis()
        if r is None:
            return None
        return r.get(_att_text_key(session_id, attachment_id))
    except Exception:
        return None


def _redis_set_text(session_id: str, attachment_id: str, text: str) -> None:
    """Store extracted text in Redis and add attachment_id to the session index."""
    try:
        r = _get_redis()
        if r is None:
            return
        ttl = _attachment_cache_ttl()
        key = _att_text_key(session_id, attachment_id)
        idx_key = _att_index_key(session_id)
        pipe = r.pipeline(transaction=False)
        pipe.setex(key, ttl, text)
        pipe.sadd(idx_key, attachment_id)
        pipe.expire(idx_key, ttl)
        pipe.execute()
    except Exception:
        logger.warning(
            "attachments.redis_write_failed",
            extra={"session_id": session_id, "attachment_id": attachment_id},
        )


def purge_session_attachment_cache(session_id: str) -> None:
    """Delete all Redis text-cache keys for a session.

    Called by :func:`delete_session_attachments` so Redis is cleaned up in
    sync with blob and MongoDB deletion.  Silent on Redis errors.
    """
    try:
        r = _get_redis()
        if r is None:
            return
        idx_key = _att_index_key(session_id)
        attachment_ids = r.smembers(idx_key) or set()
        keys_to_delete = [
            _att_text_key(session_id, aid) for aid in attachment_ids
        ] + [idx_key]
        if keys_to_delete:
            r.delete(*keys_to_delete)
        logger.info(
            "attachments.cache_purged",
            extra={"session_id": session_id, "key_count": len(keys_to_delete)},
        )
    except Exception:
        logger.warning(
            "attachments.cache_purge_failed",
            extra={"session_id": session_id},
        )

_ALLOWED_EXTENSIONS = {
    "png", "jpg", "jpeg", "webp", "gif", "bmp", "svg", "heic", "heif", "tif", "tiff",
    "pdf", "txt", "md", "csv", "json", "doc", "docx", "ppt", "pptx", "xls", "xlsx",
}
_IMAGE_EXTENSIONS = {
    "png", "jpg", "jpeg", "webp", "gif", "bmp", "svg", "heic", "heif", "tif", "tiff",
}


def _clean_filename(name: str) -> str:
    candidate = (name or "file").strip() or "file"
    candidate = re.sub(r"[^A-Za-z0-9._-]+", "_", candidate)
    return candidate[:180]


def _file_ext(name: str) -> str:
    parts = (name or "").lower().rsplit(".", 1)
    return parts[1] if len(parts) == 2 else ""


def _extract_text_for_extension(ext: str, raw: bytes) -> str:
    ext = (ext or "").lower()
    try:
        if ext in {"txt", "md"}:
            return raw.decode("utf-8", errors="replace")
        if ext == "json":
            obj = json.loads(raw.decode("utf-8", errors="replace"))
            return json.dumps(obj, indent=2, ensure_ascii=False)
        if ext == "csv":
            rows = []
            reader = csv.reader(io.StringIO(raw.decode("utf-8", errors="replace")))
            for idx, row in enumerate(reader):
                rows.append(", ".join(row))
                if idx >= 200:
                    break
            return "\n".join(rows)
        if ext == "pdf":
            try:
                from pypdf import PdfReader
            except Exception:
                return ""
            reader = PdfReader(io.BytesIO(raw))
            out = []
            for page in reader.pages[:50]:
                out.append(page.extract_text() or "")
            return "\n".join(out)
        if ext == "docx":
            try:
                from docx import Document
            except Exception:
                return ""
            doc = Document(io.BytesIO(raw))
            return "\n".join(p.text for p in doc.paragraphs)
        if ext == "pptx":
            try:
                from pptx import Presentation
            except Exception:
                return ""
            prs = Presentation(io.BytesIO(raw))
            out = []
            for slide in prs.slides[:50]:
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text:
                        out.append(text)
            return "\n".join(out)
        if ext == "xlsx":
            try:
                import openpyxl
            except Exception:
                return ""
            wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            out = []
            for sheet in wb.worksheets:
                out.append(f"=== Sheet: {sheet.title} ===")
                for row in sheet.iter_rows(values_only=True):
                    row_parts = [str(cell) if cell is not None else "" for cell in row]
                    if any(p.strip() for p in row_parts):
                        out.append("\t".join(row_parts))
            wb.close()
            return "\n".join(out)
        if ext == "xls":
            try:
                import xlrd
            except Exception:
                return ""
            wb = xlrd.open_workbook(file_contents=raw)
            out = []
            for sheet in wb.sheets():
                out.append(f"=== Sheet: {sheet.name} ===")
                for row_idx in range(sheet.nrows):
                    row_parts = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                    if any(p.strip() for p in row_parts):
                        out.append("\t".join(row_parts))
            return "\n".join(out)
    except Exception:
        logger.exception("attachments.extract_failed", extra={"extension": ext})
        return ""
    return ""


def _validate_files(files: list[UploadedFile]) -> None:
    if not files:
        raise ValueError("At least one file is required.")
    if len(files) > _MAX_ATTACHMENTS_PER_MESSAGE:
        raise ValueError(f"A maximum of {_MAX_ATTACHMENTS_PER_MESSAGE} files can be uploaded at once.")

    for item in files:
        ext = _file_ext(item.name)
        if ext not in _ALLOWED_EXTENSIONS:
            raise ValueError(f"Unsupported file type '{ext or 'unknown'}' for '{item.name}'.")
        if (item.size or 0) > _MAX_ATTACHMENT_BYTES:
            raise ValueError(f"File '{item.name}' exceeds 20 MB limit.")


def _build_blob_key(*, session_id: str, attachment_id: str, filename: str) -> str:
    """Stable, permanent blob key for this attachment.

    Format: ``sessions/{session_id}/attachments/{attachment_id}/{filename}``

    The key has no message/discussion segment because those IDs are unknown at
    upload time (attachments are staged before a message is persisted).
    ``get_attachment_content`` looks up the key from MongoDB, so the path shape
    only matters for Azure prefix-based cleanup (``delete_session_attachments``
    uses the ``sessions/{session_id}/`` prefix which still matches).
    """
    return f"sessions/{session_id}/attachments/{attachment_id}/{filename}"


def _attachment_descriptor(doc: dict) -> dict:
    return {
        "id": doc.get("attachment_id", ""),
        "filename": doc.get("filename", ""),
        "mime_type": doc.get("mime_type", "application/octet-stream"),
        "size_bytes": int(doc.get("size_bytes") or 0),
        "is_image": bool(doc.get("is_image", False)),
        "extension": doc.get("extension", ""),
    }


@traced_function("service.attachments.upload")
def upload_session_attachments(*, session: dict, files: list[UploadedFile]) -> list[dict]:
    _validate_files(files)

    strategy = build_storage_strategy()
    col = get_collection(ATTACHMENTS_COLLECTION)
    session_id = session.get("session_id", "")
    project_id = session.get("project_id", "")
    placeholders = []

    for uploaded in files:
        filename = _clean_filename(uploaded.name)
        ext = _file_ext(filename)
        attachment_id = str(uuid4())
        key = _build_blob_key(
            session_id=session_id,
            attachment_id=attachment_id,
            filename=filename,
        )

        raw = uploaded.read()
        strategy.upload_bytes(
            key=key,
            data=raw,
            content_type=(uploaded.content_type or "application/octet-stream"),
        )

        doc = {
            "attachment_id": attachment_id,
            "project_id": project_id,
            "session_id": session_id,
            "message_id": None,
            "filename": filename,
            "extension": ext,
            "mime_type": uploaded.content_type or "application/octet-stream",
            "size_bytes": int(uploaded.size or len(raw)),
            "is_image": ext in _IMAGE_EXTENSIONS,
            "blob_key": key,
            "uploaded_at": util.utc_now(),
        }
        col.insert_one(doc)
        placeholders.append(_attachment_descriptor(doc))

    logger.info(
        "attachments.uploaded",
        extra={"session_id": session_id, "count": len(placeholders)},
    )
    return placeholders


def _get_attachment_docs_for_session(session_id: str, attachment_ids: Iterable[str]) -> list[dict]:
    wanted = [str(x).strip() for x in (attachment_ids or []) if str(x).strip()]
    if not wanted:
        return []
    col = get_collection(ATTACHMENTS_COLLECTION)
    cursor = col.find({"session_id": session_id, "attachment_id": {"$in": wanted}})
    docs = list(cursor)
    index = {d.get("attachment_id"): d for d in docs}
    return [index[aid] for aid in wanted if aid in index]


@traced_function("service.attachments.bind_to_message")
def bind_attachments_to_message(*, session_id: str, message_id: str, attachment_ids: Iterable[str]) -> list[dict]:
    docs = _get_attachment_docs_for_session(session_id, attachment_ids)
    if not docs:
        return []

    col = get_collection(ATTACHMENTS_COLLECTION)
    attachment_ids_clean = [d["attachment_id"] for d in docs]
    col.update_many(
        {"session_id": session_id, "attachment_id": {"$in": attachment_ids_clean}},
        {"$set": {"message_id": message_id, "bound_at": util.utc_now()}},
    )

    return [
        {
            "id": d.get("attachment_id", ""),
            "filename": d.get("filename", ""),
            "mime_type": d.get("mime_type", "application/octet-stream"),
            "size_bytes": int(d.get("size_bytes") or 0),
            "is_image": bool(d.get("is_image", False)),
            "extension": d.get("extension", ""),
            "content_url": f"/chat/sessions/{session_id}/attachments/{d.get('attachment_id', '')}/content/",
        }
        for d in docs
    ]


def build_attachment_context_block(*, session_id: str, attachment_ids: Iterable[str]) -> str:
    """Build a text block describing non-image attachments for the agent task.

    Text extraction is lazy-cached in Redis:

    * **Cache hit** — returns immediately (< 1 ms).
    * **Cache miss + success** — downloads from Azure Blob, extracts, stores in
      Redis with ``REDIS_ATTACHMENT_TTL_SECONDS`` TTL (default 24 h), then
      returns.  Genuinely empty documents (e.g. scanned PDFs with no text
      layer) are cached so repeated runs avoid repeated blob downloads.
    * **Cache miss + exception** — blob download or extraction error: the
      result is **not** written to Redis so the next run retries the
      extraction.  Only durable successes (including empty text) are cached.

    The full extracted text is sent to the agent (no truncation) so the agent
    can see the complete document.  Redis memory use is bounded by the TTL and
    by ``REDIS_ATTACHMENT_TTL_SECONDS`` env var (default 24 h).

    Images are excluded here; they are handled separately by
    :func:`load_images_for_agents` as ``MultiModalMessage`` content.
    """
    docs = _get_attachment_docs_for_session(session_id, attachment_ids)
    if not docs:
        return ""

    strategy = build_storage_strategy()
    lines = ["", "---", "Attachments:"]
    has_content = False

    for d in docs:
        if d.get("is_image"):
            # Images are passed as pixel data in MultiModalMessage, not text.
            lines.append(f"- [image] {d.get('filename', 'image')} ({d.get('mime_type', 'image/*')})")
            has_content = True
            continue

        filename = d.get("filename", "file")
        ext = (d.get("extension") or "").lower()
        attachment_id = d.get("attachment_id", "")
        size_bytes = int(d.get("size_bytes") or 0)

        # 1. Try Redis cache.
        text = _redis_get_text(session_id, attachment_id)
        cache_hit = text is not None
        extract_failed = False

        if not cache_hit:
            # 2. Download from blob and extract.
            try:
                raw = strategy.download_bytes(key=d.get("blob_key", ""))
                text = _extract_text_for_extension(ext, raw)
            except Exception:
                logger.exception(
                    "attachments.extract_on_run_failed",
                    extra={"session_id": session_id, "attachment_id": attachment_id},
                )
                text = ""
                extract_failed = True

            # 3. Populate cache only when extraction did not raise.
            #    Exception-caused empties are NOT cached so the next run
            #    retries the blob download / extraction (e.g. transient Azure
            #    Blob timeout or pypdf error).  Genuinely empty documents are
            #    cached normally to avoid repeated blob downloads.
            if attachment_id and not extract_failed:
                _redis_set_text(session_id, attachment_id, text or "")

        logger.debug(
            "attachments.context_block",
            extra={
                "session_id": session_id,
                "attachment_id": attachment_id,
                "cache_hit": cache_hit,
                "extract_failed": extract_failed,
                "text_chars": len(text) if text else 0,
            },
        )

        if not (text and text.strip()):
            logger.warning(
                "attachments.extract_empty",
                extra={
                    "session_id": session_id,
                    "attachment_id": attachment_id,
                    "filename": filename,
                    "ext": ext,
                    "cache_hit": cache_hit,
                    "extract_failed": extract_failed,
                },
            )

        lines.append(f"\n### {filename} ({ext.upper()}, {size_bytes} bytes)")
        if text and text.strip():
            lines.append(text.strip())
        else:
            lines.append("(no extractable text content)")
        has_content = True

    if not has_content:
        return ""
    return "\n".join(lines)


def load_images_for_agents(
    *, session_id: str, attachment_ids: Iterable[str]
) -> list[tuple[str, bytes, str]]:
    """Return ``(filename, raw_bytes, mime_type)`` for each image attachment.

    Images are downloaded directly from Azure Blob on every run/resume — they
    are NOT cached in Redis (raw bytes are too large and Redis is not a blob
    store).  Each download is a single Azure GET (~50–200 ms).

    Non-image attachments are skipped; their text is handled by
    :func:`build_attachment_context_block`.

    Download failures are logged and skipped individually so a bad blob never
    aborts the whole agent run.
    """
    docs = _get_attachment_docs_for_session(session_id, attachment_ids)
    if not docs:
        return []

    strategy = build_storage_strategy()
    results: list[tuple[str, bytes, str]] = []
    for d in docs:
        if not d.get("is_image"):
            continue
        try:
            raw = strategy.download_bytes(key=d.get("blob_key", ""))
            results.append((
                d.get("filename", "image"),
                raw,
                d.get("mime_type", "image/png"),
            ))
        except Exception:
            logger.exception(
                "attachments.load_image_failed",
                extra={
                    "session_id": session_id,
                    "attachment_id": d.get("attachment_id", ""),
                    "filename": d.get("filename", ""),
                },
            )
    return results


def get_attachment_content(*, session_id: str, attachment_id: str) -> tuple[bytes, str, str]:
    col = get_collection(ATTACHMENTS_COLLECTION)
    doc = col.find_one({"session_id": session_id, "attachment_id": attachment_id})
    if not doc:
        raise ValueError("Attachment not found.")

    strategy = build_storage_strategy()
    raw = strategy.download_bytes(key=doc.get("blob_key", ""))
    return raw, doc.get("mime_type", "application/octet-stream"), doc.get("filename", "attachment")


@traced_function("service.attachments.delete_session")
def delete_session_attachments(session_id: str) -> None:
    if not session_id:
        return
    try:
        ObjectId(session_id)
    except (InvalidId, TypeError):
        return

    col = get_collection(ATTACHMENTS_COLLECTION)

    # 1. Purge Redis text cache first (uses the index set populated at run time).
    purge_session_attachment_cache(session_id)

    # 2. Delete all blob objects under the session prefix.  Always attempt this
    #    regardless of whether metadata rows exist — blobs and metadata can fall
    #    out of sync (e.g. upload interrupted after blob write, or metadata rows
    #    already removed), and skipping the prefix delete would orphan blobs.
    deleted_blobs = 0
    try:
        strategy = build_storage_strategy()
        prefix = f"sessions/{session_id}/"
        deleted_blobs = strategy.delete_prefix(prefix=prefix)
    except Exception:
        logger.exception(
            "attachments.blob_delete_failed",
            extra={"session_id": session_id},
        )

    # 3. Delete MongoDB metadata (no-op if no rows exist).
    result = col.delete_many({"session_id": session_id})

    logger.info(
        "attachments.session_deleted",
        extra={
            "session_id": session_id,
            "deleted_blob_count": int(deleted_blobs),
            "deleted_metadata_count": int(result.deleted_count),
        },
    )
